import os
from datetime import date

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


# Delete the powerpoint if it already exists
if os.path.exists(r"C:\Users\Matth\git\DataAnalysisWorkbooks\Covid19\Covid19_dashboard.pptx"):
    os.remove(r"C:\Users\Matth\git\DataAnalysisWorkbooks\Covid19\Covid19_dashboard.pptx")
    
    
# Initialize a presentation object (Note: close all open instances of powerpoint)
prs = Presentation()

# Set slide dimensions to widescreen
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Other configs
slide_numbers = True
rice_logo = True


# Useful functions

def slide_title(l, t, w, h, text, textFont, textSize, alignment):  # Text for title slide
    
    left = Inches(l)
    top = Inches(t)
    width = Inches(w)
    height = Inches(h)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tb = text_box.text_frame
    tb.word_wrap = True
    
    p = tb.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    
    font = run.font
    font.name = textFont
    font.size = Pt(textSize)
    
    return None


def slide_image(imgType, img_name, l, t, w, slideTitle, caption):
    
    if imgType == 'scatter': path = 'C:/Users/Matth/git/DataAnalysisWorkbooks/Covid19/Figures/scatter_plots/' + img_name + '.png'
        
    if imgType == 'box': path = 'C:/Users/Matth/git/DataAnalysisWorkbooks/Covid19/Figures/box_plots/' + img_name + '.png'
    
    left = l
    top = t
    img = slide.shapes.add_picture(path, left, top, w)
    
    slide_title(1.7, 0.6, 14, 3, slideTitle, 'Arial', 35, PP_ALIGN.CENTER)  # Title
    slide_title(1.6, 8, 13, 5, caption, 'Arial', 30, PP_ALIGN.CENTER)
    
    return None


# Title slide

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

slide_title(2.3, 2.45, 11, 3, 'Covid-19 Daily Statistics Dashboard', 'Arial', 60, PP_ALIGN.CENTER)  # Title
slide_title(3.5, 5, 8.4, 3, 'Matthew Decaro, Rice University', 'Arial', 30, PP_ALIGN.CENTER)  # Subtitle (name)
slide_title(3.5, 7.3, 8.4, 3, date.today().strftime('%B %#d, %Y'), 'Arial', 30, PP_ALIGN.CENTER) # Date


# Create a basic image slide with a caption

# Image slides (Note: blade_slide_layout cannot be called in a function. Must manually call it each time a new slide is made.)

continents = ['Africa', 'Asia', 'Europe', 'NorthAmerica', 'SouthAmerica', 'Oceania']
cont_caption = ['Africa', 'Asia', 'Europe', 'North America', 'South America', 'Oceania']

countries = ['UnitedStates', 'Canada', 'Mexico', 'UnitedKingdom', 'France', 'Germany', 'Japan', 'SouthKorea']
country_caption = ['United States', 'Canada', 'Mexico', 'United Kingdom', 'France', 'Germany', 'Japan', 'South Korea']

# Text slide describing Covid-19 dataset

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

slide_title(2.3, 0.6, 11, 3, 'Covid-19 Dataset', 'Arial', 40, PP_ALIGN.CENTER)  # Title
slide_title(0.7, 2.2, 20, 20, "• Data collected by 'Our World in Data' (owid).", 'Arial', 30, PP_ALIGN.LEFT)  
slide_title(0.7, 3.0, 20, 20, "• Scraped from several sources: ", 'Arial', 30, PP_ALIGN.LEFT)  
slide_title(1.7, 3.8, 20, 20, "• Specialized institutions (WHO, JHU, ...)", 'Arial', 30, PP_ALIGN.LEFT)  
slide_title(1.7, 4.6, 20, 20, "• International agencies (UN, World Bank, ...)", 'Arial', 30, PP_ALIGN.LEFT)
slide_title(1.7, 5.4, 20, 20, "• Government sources (US Department of Health & Human Services)", 'Arial', 30, PP_ALIGN.LEFT)
slide_title(0.7, 6.2, 20, 20, "• >150K rows of data, updated daily.", 'Arial', 30, PP_ALIGN.LEFT)  
slide_title(0.7, 7.0, 20, 20, "• Confirmed Covid cases, deaths, vaccinations and more for >200 countries. ", 'Arial', 30, PP_ALIGN.LEFT)  


# Total Covid-19 cases by continent
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('scatter', f'total_cases', Inches(-0.8), Inches(2), Inches(17.5), f'Total Covid-19 Cases by Continent', '')

# Total Covid-19 deaths by continent
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('scatter', f'total_deaths', Inches(-0.8), Inches(2), Inches(17.5), f'Total Covid-19 Deaths by Continent', '')

# Cases per capita per continent
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('scatter', f'casesPerCapitaPerContinent', Inches(-0.8), Inches(2), Inches(17.5), f'Covid-19 cases per capita per continent', '')

# Deaths per capita per continent
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('scatter', f'deathsPerCapitaPerContinent', Inches(-0.8), Inches(2), Inches(17.5), f'Covid-19 deaths per capita per continent', '')

# Death-case ratio per continent
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('scatter', f'deathsCaseRatio', Inches(-0.8), Inches(2), Inches(17.5), f'Total Covid-19 Deaths/Cases per continent', '')

# Global daily Covid cases
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('scatter', f'newCases_World', Inches(-0.8), Inches(2), Inches(17.5), f'Daily Global Covid-19 Cases', '')

# Global daily Covid deaths
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('scatter', f'newDeaths_World', Inches(-0.8), Inches(2), Inches(17.5), f'Daily Global Covid-19 Deaths', '')

#Box-and-whisker plots
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_image('box', f'casesPerCapita', Inches(0.5), Inches(2), Inches(5), f'Cases, Deaths, and Vaccinations per capita per country', '')
slide_image('box', f'deathsPerCapita', Inches(5.5), Inches(2), Inches(5), f'', '')
slide_image('box', f'vaccinatedPerCapita', Inches(10.5), Inches(2), Inches(5), f'', '')

# Title slide for next image slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_title(3.3, 2.45, 9, 3, 'Daily Covid-19 Cases by Continent', 'Arial', 60, PP_ALIGN.CENTER)  # Title
slide_title(3.5, 5, 8.4, 3, '(Africa, Asia, Europe, North America, South America, Oceania)', 'Arial', 30, PP_ALIGN.CENTER)
    
# Daily Covid Cases by continent
for i, continent in enumerate(continents):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_image('scatter', f'newCases_{continent}', Inches(-0.8), Inches(2), Inches(17.5), f'Daily Covid-19 Cases in {cont_caption[i]}', '')

# Title slide for next image slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide_title(3.3, 2.45, 9, 3, 'Daily Covid-19 Deaths by Continent', 'Arial', 60, PP_ALIGN.CENTER)  # Title
slide_title(3.5, 5, 8.4, 3, '(Africa, Asia, Europe, North America, South America, Oceania)', 'Arial', 30, PP_ALIGN.CENTER)    

# Daily Covid Deaths by continent
for i, continent in enumerate(continents):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_image('scatter', f'newDeaths_{continent}', Inches(-0.8), Inches(2), Inches(17.5), f'Daily Covid-19 Deaths in {cont_caption[i]}', '')

    

# Add slide numbers to footers

if slide_numbers:
    slides = prs.slides

    # Hand-picked positions based on choice of slide size (widescreen) and fontsize (27)
    left = Inches(15.1)
    top = Inches(8.1)
    width = Inches(1)
    height = Inches(1)

    for slide in slides:
        if slides.index(slide) == 0: continue  # Skip labelling the title slide number as '0'

        # Initialize the text box. Set the text, font, and fontsize
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tb = text_box.text_frame
        
        p = tb.paragraphs[0]
        run = p.add_run()
        run.text = str(slides.index(slide))
        
        font = run.font
        font.name = 'Arial'
        font.size = Pt(27)
        
        #TODO: Add a toggle for other slide number positions (bottom left, bottom center, etc.)


skip_slides = []  # Add slides to skip logo here.

if rice_logo:
    slides = prs.slides

    # Hand-picked positions based on choice of slide size (widescreen) and fontsize (27)
    left = Inches(0.3)
    top = Inches(0.2)
    path = r"C:\Users\Matth\Documents\saved_images\rice_logo.png"

    for index, slide in enumerate(slides):
        if index in skip_slides: continue
        img = slide.shapes.add_picture(path, left, top)
        
        #TODO: Add a toggle for other logo positions (top-right, bottom-left, etc.)        
        
prs.save(r"C:\Users\Matth\git\DataAnalysisWorkbooks\Covid19\Covid19_dashboard.pptx")  # Save the file